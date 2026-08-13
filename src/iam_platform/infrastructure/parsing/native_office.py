"""Word, PowerPoint and Excel read directly, without the layout models.

**Why not just use docling for these.** Docling infers structure from a
*rendered image* of a page, which is the right tool when structure is only
visible visually — a scan, a PDF with no text layer. An OOXML file is not that
document: it already states, in XML, that this run is Heading 2, that this is
a table cell, that this shape is on slide 4. Reconstructing that from pixels
is slower by orders of magnitude and *loses* information that was sitting
there in the markup. The same argument this codebase already made for
text-layer PDFs in `fast_pdf.py`, applied to the formats where it is even
clearer.

**What that buys, beyond speed.** Structure survives into `source_location`,
so a citation can say "Slide 7" or "Sheet1 rows 41-80" instead of naming the
file. Docling's markdown output flattens all of that into headings we then
have to guess at.

These parsers deliberately do **not** OCR embedded pictures. A deck that is
mostly screenshots gets little text here, and the dispatcher notices and
hands it to docling — see `dispatcher.py`. Trying to do both jobs in one
parser is how you end up with a slow path that is also bad at its fast case.
"""

from __future__ import annotations

import asyncio
from io import BytesIO
from typing import Any

from iam_platform.application.ai_resources.exceptions import DocumentParseError
from iam_platform.application.ai_resources.ports import ParsedBlock

#: Rows batched into one block. Matches the CSV parser's grouping: one row per
#: block would defeat the chunker (a 5,000-row sheet becoming 5,000 chunks),
#: and a whole sheet in one block would make "which row?" unanswerable.
_ROWS_PER_BLOCK = 40


def _clean(text: str) -> str:
    return " ".join(text.split())


def _markdown_table(rows: list[list[str]]) -> str:
    """A table as markdown, because that is what the embedding model reads.

    Pipe-delimited rows keep column association in the text itself, so a chunk
    containing "| Refund window | 30 days |" still answers a question about
    refund windows. A bare space-joined dump loses which value belongs to
    which heading.
    """
    cleaned = [[_clean(str(cell)) for cell in row] for row in rows]
    cleaned = [row for row in cleaned if any(cell for cell in row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    padded = [row + [""] * (width - len(row)) for row in cleaned]
    header, *body = padded
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


class NativeDocxParser:
    """Word, via python-docx, preserving the heading path.

    `iter_inner_content()` yields paragraphs and tables in document order --
    iterating `doc.paragraphs` then `doc.tables` separately, as most examples
    do, silently reorders every table to the end of the document.
    """

    async def parse(self, *, data: bytes, filename: str) -> list[ParsedBlock]:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> list[ParsedBlock]:
        try:
            from docx import Document
            from docx.table import Table
        except ImportError as exc:  # pragma: no cover - dependency is declared
            raise DocumentParseError(f"{filename}: Word support is unavailable") from exc

        try:
            document = Document(BytesIO(data))
        except Exception as exc:
            raise DocumentParseError(f"{filename}: could not be read as a Word document ({exc})") from exc

        blocks: list[ParsedBlock] = []
        headings: list[str] = []

        for item in document.iter_inner_content():
            if isinstance(item, Table):
                rows = [[cell.text for cell in row.cells] for row in item.rows]
                table = _markdown_table(rows)
                if table:
                    blocks.append(ParsedBlock(text=table, source_location=_where(headings, "table")))
                continue

            text = _clean(item.text)
            if not text:
                continue

            style = (item.style.name if item.style else "") or ""
            lowered = style.lower()
            if lowered == "title":
                headings = [text]
            elif lowered.startswith("heading"):
                level = _heading_level(style)
                headings = headings[: max(0, level - 1)] + [text]
            blocks.append(ParsedBlock(text=text, source_location=_where(headings, None)))

        return blocks

    def image_count(self, data: bytes) -> int:
        """How many inline pictures the document holds.

        Read by the dispatcher: a Word file with little text *and* pictures is
        probably a scan pasted into Word, which docling can read and this
        parser cannot.
        """
        try:
            from docx import Document

            return len(Document(BytesIO(data)).inline_shapes)
        except Exception:
            return 0


def _heading_level(style_name: str) -> int:
    tail = style_name.split()[-1] if style_name.split() else ""
    try:
        return int(tail)
    except ValueError:
        return 1


def _where(headings: list[str], suffix: str | None) -> str | None:
    """The heading trail as a citation, e.g. "Refunds > Eligibility"."""
    trail = " > ".join(h for h in headings if h)
    if suffix:
        return f"{trail} ({suffix})" if trail else suffix
    return trail or None


class NativePptxParser:
    """PowerPoint, via python-pptx, one block per shape in reading order.

    Shapes are sorted by position rather than taken in XML order: PowerPoint
    stores them in z-order, so an untouched iteration reads a deck in the
    order shapes happened to be added, which is rarely the order a person
    reads them in.
    """

    async def parse(self, *, data: bytes, filename: str) -> list[ParsedBlock]:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> list[ParsedBlock]:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - dependency is declared
            raise DocumentParseError(
                f"{filename}: PowerPoint support is unavailable"
            ) from exc

        try:
            presentation = Presentation(BytesIO(data))
        except Exception as exc:
            raise DocumentParseError(
                f"{filename}: could not be read as a PowerPoint file ({exc})"
            ) from exc

        blocks: list[ParsedBlock] = []
        for number, slide in enumerate(presentation.slides, start=1):
            where = f"slide {number}"
            for shape in _in_reading_order(slide.shapes):
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    table = _markdown_table(rows)
                    if table:
                        blocks.append(ParsedBlock(text=table, source_location=where))
                elif getattr(shape, "has_text_frame", False):
                    text = _clean(shape.text)
                    if text:
                        blocks.append(ParsedBlock(text=text, source_location=where))

            notes = _speaker_notes(slide)
            if notes:
                # Kept, and labelled: speaker notes routinely carry the
                # explanation the slide only gestures at, which is exactly
                # what a question tends to be about.
                blocks.append(
                    ParsedBlock(text=notes, source_location=f"{where} (speaker notes)")
                )

        return blocks

    def picture_count(self, data: bytes) -> int:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            presentation = Presentation(BytesIO(data))
            return sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
            )
        except Exception:
            return 0


def _in_reading_order(shapes: Any) -> list[Any]:
    return sorted(shapes, key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0))


def _speaker_notes(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        frame = slide.notes_slide.notes_text_frame
        return _clean(frame.text) if frame is not None else ""
    except Exception:
        return ""


class NativeXlsxParser:
    """Excel, via openpyxl in read-only mode, batched by sheet and row range.

    `read_only=True` streams rather than materialising the workbook, which is
    the difference between a 50 MB spreadsheet costing 50 MB and costing
    several gigabytes. `data_only=True` reads the cached result of a formula:
    the stored text of `=SUM(A1:A20)` is not something anyone can answer a
    question from.
    """

    async def parse(self, *, data: bytes, filename: str) -> list[ParsedBlock]:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> list[ParsedBlock]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover - dependency is declared
            raise DocumentParseError(f"{filename}: Excel support is unavailable") from exc

        try:
            workbook = load_workbook(filename=BytesIO(data), read_only=True, data_only=True)
        except Exception as exc:
            raise DocumentParseError(
                f"{filename}: could not be read as an Excel workbook ({exc})"
            ) from exc

        blocks: list[ParsedBlock] = []
        try:
            for sheet in workbook.worksheets:
                header: list[str] = []
                batch: list[list[str]] = []
                first_row = 0

                for number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    values = ["" if value is None else str(value) for value in row]
                    if not any(value.strip() for value in values):
                        continue
                    if not header:
                        # The first non-empty row is repeated at the top of
                        # every batch, so a chunk from row 900 still says what
                        # its columns mean.
                        header = values
                        continue
                    if not batch:
                        first_row = number
                    batch.append(values)
                    if len(batch) >= _ROWS_PER_BLOCK:
                        blocks.append(
                            _sheet_block(sheet.title, header, batch, first_row, number)
                        )
                        batch = []

                if batch:
                    blocks.append(
                        _sheet_block(
                            sheet.title, header, batch, first_row, first_row + len(batch) - 1
                        )
                    )
                elif header and not blocks:
                    # A sheet with a header and nothing else still says what
                    # the workbook is about.
                    blocks.append(
                        ParsedBlock(
                            text=_markdown_table([header]),
                            source_location=f"{sheet.title} (header)",
                        )
                    )
        finally:
            workbook.close()

        return blocks


def _sheet_block(
    sheet: str, header: list[str], rows: list[list[str]], first: int, last: int
) -> ParsedBlock:
    table = _markdown_table([header, *rows] if header else rows)
    where = f"{sheet} rows {first}-{last}" if last > first else f"{sheet} row {first}"
    return ParsedBlock(text=table, source_location=where)
