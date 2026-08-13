"""The extraction layer: what a file is, what it costs to open, and what
comes out.

Three properties are worth more than the rest and get the most attention here:

1. **Content decides the parser, not the filename.** Both the extension and
   the `Content-Type` are caller-controlled. A test that only uploads
   correctly-named files can never tell whether detection works or whether the
   extension happened to be right.
2. **A ZIP-based document is a decompression instruction from a stranger.**
   The bound has to be proven with an actual hostile archive, not asserted.
3. **`source_location` survives.** It is the entire reason for reading OOXML
   natively rather than rendering it -- "Slide 7" and "Sheet1 rows 41-80" are
   what a citation shows a person. A parser that returned the right text with
   no provenance would pass a naive text-equality test and be useless.

The fixtures build real files with the real libraries rather than checking in
binaries: a committed .docx is opaque, and nobody can tell from a review
whether it still contains what the test claims.
"""

from __future__ import annotations

import zipfile
from io import BytesIO

import pytest

from iam_platform.application.ai_resources.exceptions import (
    DocumentParseError,
    UnsupportedDocumentTypeError,
)
from iam_platform.core.config import IngestionSettings
from iam_platform.infrastructure.parsing.dispatcher import ParserDispatcher
from iam_platform.infrastructure.parsing.file_kind import (
    ArchiveLimits,
    FileKind,
    UnsafeArchiveError,
    assert_archive_is_safe,
    detect_kind,
)
from iam_platform.infrastructure.parsing.markup import EmlParser, HtmlParser, decode_text
from iam_platform.infrastructure.parsing.native_office import (
    NativeDocxParser,
    NativePptxParser,
    NativeXlsxParser,
)
from iam_platform.infrastructure.parsing.structured import XmlParser

pytestmark = pytest.mark.unit


# --------------------------------------------------------------------------
# Builders
# --------------------------------------------------------------------------


def _docx(paragraphs: list[tuple[str, str]], table: list[list[str]] | None = None) -> bytes:
    from docx import Document

    document = Document()
    for style, text in paragraphs:
        document.add_paragraph(text, style=style or None)
    if table:
        added = document.add_table(rows=len(table), cols=len(table[0]))
        for row_index, row in enumerate(table):
            for cell_index, value in enumerate(row):
                added.cell(row_index, cell_index).text = value
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx(slides: list[tuple[str, str, str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for title, body, notes in slides:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = title
        lower = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(2))
        lower.text_frame.text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx(sheets: dict[str, list[list[object]]]) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    for name, rows in sheets.items():
        sheet = workbook.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _zip_bomb(*, entries: int = 1, declared_size: int = 200 * 1024 * 1024) -> bytes:
    """A ZIP whose central directory declares far more than it stores.

    Built by writing highly-compressible entries, which is exactly how the real
    thing works -- there is nothing exotic about it, which is the point.
    """
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for index in range(entries):
            archive.writestr(f"entry-{index}.bin", b"\0" * declared_size)
    return buffer.getvalue()


# --------------------------------------------------------------------------
# Detection
# --------------------------------------------------------------------------


class TestDetection:
    def test_a_pdf_renamed_to_docx_is_still_detected_as_a_pdf(self) -> None:
        """The security-relevant direction. Trusting the extension here would
        hand PDF bytes to the Word reader, which fails with a confusing error
        at best."""
        assert detect_kind(b"%PDF-1.7\nstuff", "invoice.docx") is FileKind.PDF

    def test_a_docx_uploaded_as_octet_stream_is_still_a_docx(self) -> None:
        """The correctness-relevant direction, and the more common one:
        browsers send `application/octet-stream` for ordinary files all the
        time, and refusing those would be a bug visible to every tenant."""
        assert detect_kind(_docx([("", "hello")]), "notes.bin") is FileKind.DOCX

    @pytest.mark.parametrize(
        ("data", "expected"),
        [
            (b"\x89PNG\r\n\x1a\n\x00", FileKind.IMAGE),
            (b"\xff\xd8\xff\xe0junk", FileKind.IMAGE),
            (b"GIF89a....", FileKind.IMAGE),
            (b"II*\x00rest", FileKind.IMAGE),
            (b"RIFF\x00\x00\x00\x00WEBPVP8 ", FileKind.IMAGE),
            (b"{\\rtf1\\ansi", FileKind.RTF),
            (b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1rest", FileKind.LEGACY_OFFICE),
        ],
    )
    def test_signatures_are_recognised_without_any_filename_help(
        self, data: bytes, expected: FileKind
    ) -> None:
        assert detect_kind(data, "no-extension-at-all") is expected

    def test_the_three_ooxml_flavours_are_told_apart(self) -> None:
        """All three start with the same four bytes. Only the member list
        distinguishes them, which is why `_classify_zip` reads it."""
        assert detect_kind(_docx([("", "x")]), "a") is FileKind.DOCX
        assert detect_kind(_pptx([("t", "b", "")]), "a") is FileKind.PPTX
        assert detect_kind(_xlsx({"S": [["h"], ["v"]]}), "a") is FileKind.XLSX

    def test_the_extension_is_the_fallback_for_formats_with_no_signature(self) -> None:
        assert detect_kind(b"a,b\n1,2\n", "data.csv") is FileKind.CSV
        assert detect_kind(b"# Title", "readme.md") is FileKind.MARKDOWN
        assert detect_kind(b"<root/>", "feed.xml") is FileKind.XML

    def test_an_unidentifiable_file_is_unknown_rather_than_guessed(self) -> None:
        """`UNKNOWN` is what makes the upload endpoint's 415 honest. Guessing
        `TEXT` here would push every executable, archive and font into the
        pipeline to be indexed as mojibake."""
        assert detect_kind(b"MZ\x90\x00binary", "setup.exe") is FileKind.UNKNOWN

    def test_a_zip_that_is_not_a_document_is_unknown(self) -> None:
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as archive:
            archive.writestr("photos/1.jpg", b"x")
        assert detect_kind(buffer.getvalue(), "album.zip") is FileKind.UNKNOWN


# --------------------------------------------------------------------------
# Archive limits
# --------------------------------------------------------------------------


class TestArchiveLimits:
    def test_an_ordinary_office_file_passes(self) -> None:
        assert_archive_is_safe(_docx([("", "hello")]))
        assert_archive_is_safe(_pptx([("t", "b", "")]))
        assert_archive_is_safe(_xlsx({"S": [["h"], ["v"]]}))

    def test_a_single_hugely_expanding_entry_is_refused(self) -> None:
        with pytest.raises(UnsafeArchiveError) as raised:
            assert_archive_is_safe(
                _zip_bomb(), limits=ArchiveLimits(max_uncompressed_bytes=64 * 1024 * 1024)
            )
        assert "uncompressed-size limit" in str(raised.value)

    def test_the_ratio_check_catches_what_the_total_size_check_misses(self) -> None:
        """A 20 MB expansion is unremarkable as a total and absurd as a ratio.
        Both checks exist because either alone leaves a shape of archive
        through: raise the total limit for a legitimate large deck and the
        total check stops seeing bombs, and small entries expand wildly for
        entirely innocent reasons."""
        with pytest.raises(UnsafeArchiveError) as raised:
            assert_archive_is_safe(
                _zip_bomb(declared_size=20 * 1024 * 1024),
                limits=ArchiveLimits(max_uncompressed_bytes=512 * 1024 * 1024),
            )
        assert "which is not a document" in str(raised.value)

    def test_too_many_entries_is_refused_before_anything_is_read(self) -> None:
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as archive:
            for index in range(60):
                archive.writestr(f"{index}.txt", b"x")
        with pytest.raises(UnsafeArchiveError) as raised:
            assert_archive_is_safe(buffer.getvalue(), limits=ArchiveLimits(max_entries=50))
        assert "entries" in str(raised.value)

    def test_a_small_entry_with_a_wild_ratio_is_allowed(self) -> None:
        """The exemption is deliberate. A 40-byte XML stub expanding to 40 KB
        is a 1000x ratio and completely harmless; refusing it would reject
        ordinary documents in the name of a threat that needs *volume* to
        matter."""
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("tiny.xml", b"\0" * 100_000)
        assert_archive_is_safe(buffer.getvalue())

    def test_a_corrupt_archive_is_reported_as_unsafe_not_crashed_on(self) -> None:
        with pytest.raises(UnsafeArchiveError):
            assert_archive_is_safe(b"PK\x03\x04 and then nothing valid")


# --------------------------------------------------------------------------
# Native Office parsers
# --------------------------------------------------------------------------


class TestNativeDocx:
    async def test_the_heading_trail_becomes_the_source_location(self) -> None:
        blocks = await NativeDocxParser().parse(
            data=_docx(
                [
                    ("Heading 1", "Refunds"),
                    ("Heading 2", "Eligibility"),
                    ("", "Refunds are available within 30 days of purchase."),
                ]
            ),
            filename="policy.docx",
        )
        body = [b for b in blocks if b.text.startswith("Refunds are available")]
        assert body[0].source_location == "Refunds > Eligibility"

    async def test_a_sibling_heading_replaces_rather_than_nests(self) -> None:
        blocks = await NativeDocxParser().parse(
            data=_docx(
                [
                    ("Heading 1", "Refunds"),
                    ("Heading 2", "Eligibility"),
                    ("Heading 2", "Exclusions"),
                    ("", "Gift cards are excluded."),
                ]
            ),
            filename="policy.docx",
        )
        body = [b for b in blocks if b.text.startswith("Gift cards")]
        assert body[0].source_location == "Refunds > Exclusions"

    async def test_a_table_keeps_column_association_as_markdown(self) -> None:
        """A space-joined dump would put "30 days" in the chunk without
        anything tying it to "Refund window", so the passage would retrieve
        and then fail to answer."""
        blocks = await NativeDocxParser().parse(
            data=_docx([("", "intro")], table=[["Term", "Value"], ["Refund window", "30 days"]]),
            filename="terms.docx",
        )
        tables = [b for b in blocks if b.text.startswith("|")]
        assert "| Refund window | 30 days |" in tables[0].text

    async def test_a_file_that_is_not_a_word_document_fails_with_its_own_reason(self) -> None:
        with pytest.raises(DocumentParseError):
            await NativeDocxParser().parse(data=b"nonsense", filename="broken.docx")


class TestNativePptx:
    async def test_every_block_carries_its_slide_number(self) -> None:
        blocks = await NativePptxParser().parse(
            data=_pptx([("Agenda", "Introductions", ""), ("Pricing", "Tiered by seat", "")]),
            filename="deck.pptx",
        )
        by_text = {b.text: b.source_location for b in blocks}
        assert by_text["Agenda"] == "slide 1"
        assert by_text["Tiered by seat"] == "slide 2"

    async def test_speaker_notes_are_kept_and_labelled(self) -> None:
        """They routinely carry the explanation the slide only gestures at,
        which is what a question is usually about. Labelling them separately
        keeps a citation honest about where the answer came from."""
        blocks = await NativePptxParser().parse(
            data=_pptx([("Pricing", "Tiered", "Discounts are negotiable up to 15%.")]),
            filename="deck.pptx",
        )
        notes = [b for b in blocks if "negotiable" in b.text]
        assert notes[0].source_location == "slide 1 (speaker notes)"


class TestNativeXlsx:
    async def test_rows_are_batched_and_the_range_is_recorded(self) -> None:
        rows: list[list[object]] = [["id", "answer"]]
        rows += [[index, f"answer {index}"] for index in range(1, 101)]
        blocks = await NativeXlsxParser().parse(
            data=_xlsx({"FAQ": rows}), filename="faq.xlsx"
        )
        assert len(blocks) == 3
        assert blocks[0].source_location == "FAQ rows 2-41"
        assert blocks[-1].source_location == "FAQ rows 82-101"

    async def test_the_header_is_repeated_in_every_batch(self) -> None:
        """Without this, a chunk drawn from row 900 is a grid of bare values
        with nothing saying which column is which."""
        rows: list[list[object]] = [["question", "answer"]]
        rows += [[f"q{index}", f"a{index}"] for index in range(1, 61)]
        blocks = await NativeXlsxParser().parse(
            data=_xlsx({"FAQ": rows}), filename="faq.xlsx"
        )
        assert all(block.text.startswith("| question | answer |") for block in blocks)

    async def test_each_sheet_is_scoped_separately(self) -> None:
        blocks = await NativeXlsxParser().parse(
            data=_xlsx(
                {
                    "Refunds": [["term", "value"], ["window", "30 days"]],
                    "Shipping": [["term", "value"], ["carrier", "DHL"]],
                }
            ),
            filename="book.xlsx",
        )
        assert {block.source_location for block in blocks} == {
            "Refunds row 2",
            "Shipping row 2",
        }


# --------------------------------------------------------------------------
# Markup
# --------------------------------------------------------------------------


class TestMarkup:
    def test_a_declared_charset_beats_statistical_detection(self) -> None:
        """Detection needs a few hundred bytes to mean anything. On this
        twelve-byte string charset-normalizer confidently decodes `café` into
        an Arabic presentation form -- measured, not hypothesised -- and short
        documents (an HTML fragment, an exported email) are ordinary here. The
        declaration is a statement by whoever wrote the file, so it is used
        first."""
        raw = "Prix du café".encode("cp1252")
        assert "café" in decode_text(raw, declared="iso-8859-1")

    def test_a_long_windows_encoded_file_is_detected_without_a_declaration(self) -> None:
        """The path a plain `.txt` export takes, where nothing declares
        anything. It works once there is enough text for statistics."""
        raw = (
            "Le prix du café a augmenté cette année. Les clients privilégiés "
            "bénéficient d'une réduction. Veuillez consulter la politique de "
            "remboursement ci-après."
        ).encode("cp1252")
        assert "café" in decode_text(raw)

    def test_utf8_wins_over_a_stale_declaration(self) -> None:
        """A template that still says `iso-8859-1` while emitting UTF-8 is an
        everyday occurrence; bytes that decode cleanly as UTF-8 essentially
        never do so by accident. So the declaration loses this one."""
        assert decode_text("café".encode(), declared="iso-8859-1") == "café"

    def test_an_html_meta_charset_is_honoured(self) -> None:
        html = '<html><head><meta charset="windows-1252"></head><body><p>café</p></body></html>'
        assert "café" in decode_text(html.encode("cp1252"))

    def test_a_utf8_bom_is_not_indexed_as_a_character(self) -> None:
        assert decode_text("﻿Hello".encode()) == "Hello"

    async def test_scripts_and_styles_are_not_indexed_as_prose(self) -> None:
        html = b"""
        <html><head><style>.a{color:red}</style><script>var x=1;</script></head>
        <body><h1>Refunds</h1><p>Within 30 days.</p></body></html>
        """
        blocks = await HtmlParser().parse(data=html, filename="page.html")
        combined = " ".join(block.text for block in blocks)
        assert "var x" not in combined
        assert "color:red" not in combined
        assert "Within 30 days." in combined

    async def test_html_headings_become_the_source_location(self) -> None:
        html = b"<html><body><h1>Refunds</h1><h2>Eligibility</h2><p>30 days.</p></body></html>"
        blocks = await HtmlParser().parse(data=html, filename="page.html")
        body = [b for b in blocks if b.text == "30 days."]
        assert body[0].source_location == "Refunds > Eligibility"

    async def test_email_prefers_the_plain_part_over_the_html_one(self) -> None:
        """The plain part is what the sender wrote. The HTML alternative is
        the same words wrapped in a signature block and three quote levels of
        styling."""
        raw = (
            b"Subject: Refund request\r\n"
            b"From: a@example.com\r\n"
            b"To: support@example.com\r\n"
            b'Content-Type: multipart/alternative; boundary="B"\r\n\r\n'
            b"--B\r\nContent-Type: text/plain\r\n\r\nPlease refund order 42.\r\n"
            b"--B\r\nContent-Type: text/html\r\n\r\n<p>Please refund order 42.</p>"
            b"<p>Sent from my phone</p>\r\n--B--\r\n"
        )
        blocks = await EmlParser().parse(data=raw, filename="thread.eml")
        body = [b for b in blocks if b.source_location == "body"]
        assert body[0].text == "Please refund order 42."
        assert "Sent from my phone" not in " ".join(b.text for b in blocks)

    async def test_email_headers_are_indexed_so_the_thread_is_searchable(self) -> None:
        raw = b"Subject: Refund request\r\nFrom: a@example.com\r\n\r\nBody text.\r\n"
        blocks = await EmlParser().parse(data=raw, filename="thread.eml")
        assert blocks[0].source_location == "headers"
        assert "Subject: Refund request" in blocks[0].text


# --------------------------------------------------------------------------
# XML entity expansion
# --------------------------------------------------------------------------


class TestXmlEntities:
    """The same class of attack as the ZIP guards: tenant input deciding what
    the worker allocates.

    Worth its own section because the code previously carried a *comment*
    asserting this was safe. It was half right -- external entities really are
    not resolved by `ElementTree` -- and the wrong half was the one that
    mattered.
    """

    def _bomb(self, levels: int) -> bytes:
        lines = [b'<?xml version="1.0"?>', b"<!DOCTYPE lolz [", b' <!ENTITY lol "lol">']
        for level in range(1, levels + 1):
            previous = b"lol" if level == 1 else b"lol%d" % (level - 1)
            lines.append(b' <!ENTITY lol%d "%s">' % (level, (b"&" + previous + b";") * 10))
        lines += [b"]>", b"<root>&lol%d;</root>" % levels]
        return b"\n".join(lines)

    async def test_an_entity_expansion_bomb_is_refused(self) -> None:
        """Measured before the fix: four levels through the *real* parser
        produced 30,006 characters from ~400 bytes, and each further level
        multiplies by ten. The element cap cannot help -- expansion happens
        during parsing, before there is anything to count."""
        with pytest.raises(DocumentParseError) as raised:
            await XmlParser().parse(
                data=self._bomb(4), content_type="application/xml", filename="bomb.xml"
            )
        assert "entities" in str(raised.value)

    async def test_an_external_entity_reference_is_refused(self) -> None:
        """XXE. The worker can reach the filesystem and the internal network,
        so a resolved external entity is a file-read or SSRF primitive whose
        output would then be helpfully indexed and searchable."""
        xxe = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE r [<!ENTITY x SYSTEM "file:///etc/passwd">]>'
            b"<r>&x;</r>"
        )
        with pytest.raises(DocumentParseError):
            await XmlParser().parse(
                data=xxe, content_type="application/xml", filename="xxe.xml"
            )

    async def test_ordinary_xml_is_unaffected(self) -> None:
        """The half of the change that could regress silently: refusing every
        document with a DOCTYPE would be a safe-looking fix that quietly
        stopped ingesting legitimate files."""
        blocks = await XmlParser().parse(
            data=b'<catalog><item sku="A1">Refund window is 45 days</item></catalog>',
            content_type="application/xml",
            filename="catalog.xml",
        )
        assert blocks[0].text == "item: Refund window is 45 days (sku=A1)"


# --------------------------------------------------------------------------
# Routing
# --------------------------------------------------------------------------


class _NeverCalled:
    """Stands in for docling. Any call is a routing failure, so it says so
    loudly rather than recording a flag nobody asserts on."""

    def __init__(self) -> None:
        self.called = False

    async def parse(self, *, data: bytes, content_type: str, filename: str) -> list[object]:
        self.called = True
        raise AssertionError(f"docling was reached for {filename}")


class TestRouting:
    @pytest.mark.parametrize(
        ("data_factory", "filename"),
        [
            (lambda: _docx([("", "Refunds are available within 30 days.")]), "policy.docx"),
            (lambda: _pptx([("Agenda", "Introductions", "")]), "deck.pptx"),
            (lambda: _xlsx({"FAQ": [["q", "a"], ["when", "30 days"]]}), "faq.xlsx"),
            (lambda: b"a,b\n1,2\n", "data.csv"),
            (lambda: b'{"a": 1}', "data.json"),
            (lambda: b"<root><a>1</a></root>", "data.xml"),
            (lambda: b"<html><body><p>hello</p></body></html>", "page.html"),
            (lambda: b"Subject: hi\r\n\r\nbody\r\n", "mail.eml"),
            (lambda: b"# Title\n\nSome prose.", "notes.md"),
        ],
    )
    async def test_declared_structure_never_reaches_the_layout_models(
        self, data_factory: object, filename: str
    ) -> None:
        """The central claim of this package. Docling infers structure from a
        rendered page; every format here *states* its structure, so rendering
        it is slower and loses information that was already present."""
        docling = _NeverCalled()
        dispatcher = ParserDispatcher(docling=docling)

        blocks = await dispatcher.parse(
            data=data_factory(),  # type: ignore[operator]
            content_type="application/octet-stream",
            filename=filename,
        )

        assert blocks
        assert not docling.called

    async def test_a_thin_office_file_with_pictures_falls_through_to_ocr(self) -> None:
        """A scan pasted into Word: almost no text, and pictures. Neither
        signal alone is enough -- a one-line memo is legitimately thin, and a
        well-written report with a chart parses perfectly."""
        calls: list[str] = []

        class _Ocr:
            async def parse(self, *, data: bytes, content_type: str, filename: str) -> list[object]:
                calls.append(filename)
                from iam_platform.application.ai_resources.ports import ParsedBlock

                return [ParsedBlock(text="text recovered by OCR", source_location="page 1")]

        dispatcher = ParserDispatcher(
            IngestionSettings(native_office_min_chars=200), docling=_Ocr()
        )
        blocks = await dispatcher.parse(
            data=_docx_with_picture(),
            content_type="application/octet-stream",
            filename="scan.docx",
        )

        assert calls == ["scan.docx"]
        assert blocks[0].text == "text recovered by OCR"

    async def test_a_thin_office_file_with_no_pictures_is_left_alone(self) -> None:
        """Both conditions are required, and this is the half that stops the
        pipeline burning minutes of OCR to reproduce one line it already has."""
        docling = _NeverCalled()
        dispatcher = ParserDispatcher(
            IngestionSettings(native_office_min_chars=200), docling=docling
        )

        blocks = await dispatcher.parse(
            data=_docx([("", "Closed on Mondays.")]),
            content_type="application/octet-stream",
            filename="memo.docx",
        )

        assert blocks[0].text == "Closed on Mondays."
        assert not docling.called

    async def test_a_hostile_archive_is_refused_before_a_parser_opens_it(self) -> None:
        docling = _NeverCalled()
        dispatcher = ParserDispatcher(
            IngestionSettings(archive_max_uncompressed_bytes=1024), docling=docling
        )

        bomb = _bomb_shaped_like_a_docx()
        with pytest.raises(DocumentParseError):
            await dispatcher.parse(
                data=bomb,
                content_type=(
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ),
                filename="invoice.docx",
            )
        assert not docling.called

    async def test_an_executable_named_as_a_document_is_refused_on_content(self) -> None:
        """`supports()` passes it -- the extension is fine -- so the refusal
        has to come from the bytes. This is the case the extension check
        cannot see."""
        dispatcher = ParserDispatcher(docling=_NeverCalled())
        assert dispatcher.supports(content_type="application/pdf", filename="report.pdf")

        with pytest.raises(UnsupportedDocumentTypeError):
            await dispatcher.parse(
                data=b"MZ\x90\x00\x03\x00\x00\x00binary payload",
                content_type="application/pdf",
                filename="report.pdf",
            )


def _docx_with_picture() -> bytes:
    """A Word file with one line of text and an embedded image -- the shape a
    scanned page pasted into Word actually has."""
    from docx import Document
    from docx.shared import Inches

    document = Document()
    document.add_paragraph("Scan")
    document.add_picture(BytesIO(_png()), width=Inches(1))
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _png() -> bytes:
    """A 1x1 PNG, written out rather than base64-decoded so a reader can see
    that it is exactly what it claims to be."""
    import struct
    import zlib

    def chunk(tag: bytes, payload: bytes) -> bytes:
        return (
            struct.pack(">I", len(payload))
            + tag
            + payload
            + struct.pack(">I", zlib.crc32(tag + payload) & 0xFFFFFFFF)
        )

    header = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", header)
        + chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
        + chunk(b"IEND", b"")
    )


def _bomb_shaped_like_a_docx() -> bytes:
    """A **fully valid** Word document carrying one oversized extra entry.

    Validity is the whole point, and the first version of this fixture got it
    wrong: a stub archive with a fake `word/document.xml` is refused by the
    guard *and* by python-docx, so the test passed with the guard deleted. It
    proved nothing. Mutation testing caught it -- the same failure mode Phase
    12's SSRF test had.

    Built from a real document so that, without the size guard, the parse
    succeeds. The only thing that can fail this file is the guard.
    """
    original = _docx([("", "Perfectly ordinary invoice text.")])

    buffer = BytesIO()
    with (
        zipfile.ZipFile(BytesIO(original)) as source,
        zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive,
    ):
        for info in source.infolist():
            archive.writestr(info, source.read(info.filename))
        archive.writestr("customXml/payload.bin", b"\0" * (4 * 1024 * 1024))
    return buffer.getvalue()
